# loader.py
import os
from typing import Tuple

import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation
import markdown


def load_text(path: str) -> list[dict]:
    """
    파일 경로 → [{'page': 1, 'content': '...'}, ...]
    """
    ext = os.path.splitext(path)[1].lstrip(".").lower()
    results = []

    if ext == "txt":
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            content = f.read()
        results.append({"page": 1, "content": content})

    elif ext == "pdf":
        try:
            with pdfplumber.open(path) as pdf:
                for i, page in enumerate(pdf.pages):
                    text = page.extract_text() or ""
                    if text.strip():
                        results.append({"page": i + 1, "content": text})
        except Exception as e:
            print(f"Error reading PDF {path}: {e}")

    elif ext == "docx":
        doc = DocxDocument(path)
        # DOCX는 페이지 개념이 명확하지 않으므로 전체를 1페이지로 취급하거나 단락별로 나눌 수 있음
        # 여기서는 전체를 1페이지로 처리
        content = "\n".join(p.text for p in doc.paragraphs)
        results.append({"page": 1, "content": content})

    elif ext == "pptx":
        pres = Presentation(path)
        for i, slide in enumerate(pres.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
            content = "\n".join(texts)
            if content.strip():
                results.append({"page": i + 1, "content": content})

    elif ext == "md":
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            md_text = f.read()
        # Markdown도 1페이지로 처리
        results.append({"page": 1, "content": md_text})
    
    return results
