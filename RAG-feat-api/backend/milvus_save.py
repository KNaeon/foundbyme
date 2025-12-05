"""
Phase 1: ë¬¸ì„œ ë¡œë”© ë° ì¸ë±ì‹± (txtai â†’ Milvus - ë²¡í„° ì²˜ë¦¬ ìˆ˜ì •)
- ì—­í• : ./data í´ë”ì˜ ëª¨ë“  ë¬¸ì„œë¥¼ txtaië¡œ ë²¡í„°í™” â†’ Milvusì— ì €ì¥
- ìˆ˜ì •: numpy array ì²˜ë¦¬ ë° ë²¡í„° ê²€ì¦ ê°•í™” 
"""

import os
import glob
from pymilvus import connections, FieldSchema, CollectionSchema, DataType, Collection, utility
from txtai.embeddings import Embeddings
import numpy as np

# =============================================================================
# ì„¤ì •
# =============================================================================
DATA_DIR = "./data" 
MILVUS_HOST = "127.0.0.1"
MILVUS_PORT = "19530"
COLLECTION_NAME = "study_docs"

# txtai ì„¤ì • (ë²¡í„° ìƒì„± ì—”ì§„)
EMBEDDINGS_CONFIG = {
    "path": "sentence-transformers/all-MiniLM-L6-v2",
    "content": True,
}
EMBED_DIM = 384

# =============================================================================
# 1. txtai Embeddings ì´ˆê¸°í™”
# =============================================================================
def initialize_embeddings():
    """txtai Embeddings ê°ì²´ ìƒì„± (ë²¡í„° ìƒì„± ì—”ì§„)"""
    print("[TXTAI] Initializing txtai embeddings...")
    print(f"[TXTAI] Config: {EMBEDDINGS_CONFIG}")
    
    embeddings = Embeddings(EMBEDDINGS_CONFIG)
    
    print("[TXTAI] Model loaded successfully")
    print(f"[TXTAI] Embedding dimension: {EMBED_DIM}")
    return embeddings

# =============================================================================
# 2. Milvus ì—°ê²° ë° ì»¬ë ‰ì…˜ ì„¤ì •
# =============================================================================
def connect_milvus():
    """Milvus Standaloneì— ì—°ê²°"""
    print(f"\n[MILVUS] Connecting to {MILVUS_HOST}:{MILVUS_PORT}...")
    connections.connect("default", host=MILVUS_HOST, port=MILVUS_PORT)
    print("[MILVUS] Connected successfully")

def setup_milvus_collection(embed_dim: int = 384):
    """Milvus ì»¬ë ‰ì…˜ ìƒì„±"""
    
    # ê¸°ì¡´ ì»¬ë ‰ì…˜ í™•ì¸
    if COLLECTION_NAME in utility.list_collections():
        print(f"[MILVUS] Collection '{COLLECTION_NAME}' already exists. Dropping...")
        utility.drop_collection(COLLECTION_NAME)
    
    # ìŠ¤í‚¤ë§ˆ ì •ì˜
    fields = [
        FieldSchema(name="id", dtype=DataType.INT64, is_primary=True, auto_id=True),
        FieldSchema(name="vector", dtype=DataType.FLOAT_VECTOR, dim=embed_dim),
        FieldSchema(name="text", dtype=DataType.VARCHAR, max_length=65535),
        FieldSchema(name="path", dtype=DataType.VARCHAR, max_length=512),
        FieldSchema(name="filename", dtype=DataType.VARCHAR, max_length=256),
        FieldSchema(name="doc_type", dtype=DataType.VARCHAR, max_length=50),
    ]
    
    schema = CollectionSchema(fields, description="txtai â†’ milvus integration")
    collection = Collection(COLLECTION_NAME, schema)
    
    # ì¸ë±ìŠ¤ ìƒì„±
    index_params = {
        "index_type": "IVF_FLAT",
        "metric_type": "IP",
        "params": {"nlist": 1024}
    }
    collection.create_index("vector", index_params)
    collection.load()
    
    print(f"[MILVUS] Collection '{COLLECTION_NAME}' created with schema:")
    print(f"  - Embedding dimension: {embed_dim}")
    print(f"  - Index type: IVF_FLAT")
    
    return collection

# =============================================================================
# 3. íŒŒì¼ì—ì„œ í…ìŠ¤íŠ¸ ì¶”ì¶œ
# =============================================================================
def extract_text(path: str) -> str:
    """ë‹¤ì–‘í•œ í˜•ì‹ì˜ íŒŒì¼ì—ì„œ í…ìŠ¤íŠ¸ ì¶”ì¶œ"""
    lower = path.lower()
    try:
        if lower.endswith(".pdf"):
            from pypdf import PdfReader
            reader = PdfReader(path)
            return "\n".join(page.extract_text() or "" for page in reader.pages)
        
        elif lower.endswith((".docx", ".doc")):
            from docx import Document
            doc = Document(path)
            return "\n".join(p.text for p in doc.paragraphs)
        
        elif lower.endswith((".pptx", ".ppt")):
            from pptx import Presentation
            prs = Presentation(path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        texts.append(shape.text)
            return "\n".join(texts)
        
        else:  # txt ë° ê¸°íƒ€
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
    
    except Exception as e:
        print(f"[WARN] Failed to extract {path}: {e}")
        return ""

def load_all_documents(root_dir: str):
    """./data í´ë”ì˜ ëª¨ë“  ë¬¸ì„œ ë¡œë“œ"""
    print(f"\n[LOAD] Scanning {root_dir}...")
    
    documents = []
    patterns = ["**/*.pdf", "**/*.docx", "**/*.doc", "**/*.pptx", "**/*.ppt", "**/*.txt"]
    
    for pattern in patterns:
        for file_path in glob.glob(os.path.join(root_dir, pattern), recursive=True):
            text = extract_text(file_path)
            
            if text and text.strip():
                filename = os.path.basename(file_path)
                doc_type = os.path.splitext(filename)[1][1:]
                
                documents.append({
                    "path": file_path,
                    "filename": filename,
                    "text": text,
                    "doc_type": doc_type
                })
                
                print(f"  [LOAD] {filename} ({len(text)} chars)")
    
    print(f"[INFO] Total documents loaded: {len(documents)}")
    return documents

# =============================================================================
# 4. âœ… numpy array ì²˜ë¦¬ ìˆ˜ì • - txtai ë²¡í„°í™” ë° Milvus ì €ì¥
# =============================================================================
def vectorize_and_index_via_txtai(embeddings: Embeddings, collection: Collection, documents: list):
    """
    âœ… numpy array ì²˜ë¦¬ ìˆ˜ì •
    
    í”Œë¡œìš°:
    1. ë¬¸ì„œ í…ìŠ¤íŠ¸ ì¶”ì¶œ
    2. txtai.batchtransform()ìœ¼ë¡œ ë²¡í„°í™”
    3. numpy array â†’ ë¦¬ìŠ¤íŠ¸ ë³€í™˜ (ì•ˆì „í•˜ê²Œ!)
    4. ë²¡í„° â†’ Milvusì— ì €ì¥
    """
    if not documents:
        print("[ERROR] No documents to index")
        return
    
    print(f"\n[VECTORIZE] Processing {len(documents)} documents via txtai...")
    
    # âœ… Step 1: ë¬¸ì„œ í…ìŠ¤íŠ¸ ì¶”ì¶œ
    texts = [doc["text"] for doc in documents]
    print(f"[VECTORIZE] Extracted {len(texts)} texts")
    
    # âœ… Step 2: txtai.batchtransform()ìœ¼ë¡œ ë²¡í„°í™”
    print("[VECTORIZE] Vectorizing via txtai.batchtransform()...")
    
    try:
        vectors = embeddings.batchtransform(texts)
        
        print(f"[VECTORIZE] Vectorization complete via txtai")
        
        # âœ… Step 3: numpy array íƒ€ì… í™•ì¸ ë° ë³€í™˜
        print(f"[VECTORIZE] Vectors type: {type(vectors)}")
        print(f"[VECTORIZE] Vectors shape: {vectors.shape if isinstance(vectors, np.ndarray) else 'N/A'}")
        
        # numpy array í™•ì¸
        if isinstance(vectors, np.ndarray):
            print(f"[VECTORIZE] Converting numpy array to list...")
            print(f"[VECTORIZE] Array shape: {vectors.shape} (rows, cols)")
            
            # âœ… ì•ˆì „í•œ ë³€í™˜: numpy arrayë¥¼ ëª…ì‹œì ìœ¼ë¡œ ë¦¬ìŠ¤íŠ¸ë¡œ
            vectors_list = vectors.tolist()
            print(f"[VECTORIZE] Converted to list with {len(vectors_list)} vectors")
            
            # ì²« ë²¡í„° í™•ì¸
            if vectors_list and len(vectors_list) > 0:
                first_vec = vectors_list[0]
                if isinstance(first_vec, list):
                    print(f"[VECTORIZE] First vector dimension: {len(first_vec)}")
                    print(f"[VECTORIZE] First vector sample: {first_vec[:5]}...")
            
            vectors = vectors_list
        else:
            print(f"[VECTORIZE] Vectors is already list-like: {type(vectors)}")
        
        # âœ… ë²¡í„° ê°œìˆ˜ í™•ì¸ (ì•ˆì „í•˜ê²Œ - numpy array bool ë¬¸ì œ í•´ê²°)
        num_vectors = len(vectors)
        print(f"[VECTORIZE] Total vectors: {num_vectors}")
        
        # ë²¡í„°ê°€ ì—†ìœ¼ë©´ ì—ëŸ¬
        if num_vectors == 0:
            print("[ERROR] No vectors generated!")
            return
    
    except Exception as e:
        print(f"[ERROR] Vector extraction from txtai failed: {e}")
        import traceback
        traceback.print_exc()
        return
    
    # âœ… Step 4: Milvusì— ì‚½ì…
    print(f"\n[MILVUS] Preparing {len(documents)} entities for Milvus insertion...")
    
    entities = []
    for i, (doc, vector) in enumerate(zip(documents, vectors)):
        # âœ… ë²¡í„° íƒ€ì… ì¬í™•ì¸ ë° ë³€í™˜
        if isinstance(vector, np.ndarray):
            vector = vector.tolist()
        elif not isinstance(vector, list):
            vector = list(vector)
        
        # âœ… ë²¡í„° ê¸¸ì´ í™•ì¸
        if len(vector) != EMBED_DIM:
            print(f"[WARN] Vector {i} has wrong dimension: {len(vector)} (expected {EMBED_DIM})")
            continue
        
        entity = {
            "vector": vector,
            "text": doc["text"][:65000],  # Milvus VARCHAR ì œí•œ
            "path": doc["path"],
            "filename": doc["filename"],
            "doc_type": doc["doc_type"]
        }
        entities.append(entity)
    
    if not entities:
        print("[ERROR] No valid entities to insert!")
        return
    
    # âœ… Milvusì— ì¼ê´„ ì‚½ì…
    print(f"[MILVUS] Inserting {len(entities)} entities into Milvus...")
    try:
        insert_result = collection.insert(entities)
        collection.flush()
        collection.load()
        
        print(f"[MILVUS] Successfully inserted {len(insert_result.primary_keys)} documents")
        print(f"[MILVUS] Collection now has {collection.num_entities} total entities")
    
    except Exception as e:
        print(f"[ERROR] Milvus insertion failed: {e}")
        import traceback
        traceback.print_exc()

# =============================================================================
# 5. í†µê³„ ë° ê²€ì¦
# =============================================================================
def show_collection_stats(collection):
    """ì»¬ë ‰ì…˜ í†µê³„ í‘œì‹œ"""
    print(f"\n[STATS]")
    print(f"  - Collection name: {COLLECTION_NAME}")
    print(f"  - Total entities: {collection.num_entities}")
    print(f"  - Schema fields: {[field.name for field in collection.schema.fields]}")

# =============================================================================
# 6. í…ŒìŠ¤íŠ¸ ê²€ìƒ‰
# =============================================================================
def test_search(embeddings: Embeddings, collection: Collection, query: str, top_k: int = 3):
    """í…ŒìŠ¤íŠ¸ ê²€ìƒ‰ (txtaië¡œ ì¿¼ë¦¬ ë²¡í„°í™”)"""
    print(f"\n[TEST_SEARCH] Query: '{query}'")
    
    try:
        # âœ… txtai.transform()ìœ¼ë¡œ ì¿¼ë¦¬ ë²¡í„°í™”
        print("[TEST_SEARCH] Vectorizing query via txtai.transform()...")
        
        qvec = embeddings.transform(query)
        
        print(f"[TEST_SEARCH] Query vector type: {type(qvec)}")
        
        # âœ… numpy array ë³€í™˜
        if isinstance(qvec, np.ndarray):
            qvec = qvec.tolist()
            print(f"[TEST_SEARCH] Converted numpy array to list")
        
        if not isinstance(qvec, list):
            qvec = list(qvec)
        
        print(f"[TEST_SEARCH] Query vector created via txtai: dimension={len(qvec)}")
        
        # Milvusì—ì„œ ê²€ìƒ‰
        search_params = {"metric_type": "IP", "params": {"nprobe": 16}}
        results = collection.search(
            data=[qvec],
            anns_field="vector",
            param=search_params,
            limit=top_k,
            output_fields=["text", "path", "filename"]
        )
        
        print(f"[TEST_SEARCH] Top {top_k} results:")
        if results and len(results) > 0 and len(results[0]) > 0:
            for i, hit in enumerate(results[0], 1):
                entity = hit.entity
                print(f"\n  {i}. Score: {hit.score:.4f}")
                print(f"     File: {entity.get('filename')}")
                print(f"     Path: {entity.get('path')}")
                text_preview = entity.get('text', '')[:150].replace('\n', ' ')
                print(f"     Preview: {text_preview}...")
        else:
            print("  No results found (collection might be empty)")
    
    except Exception as e:
        print(f"[WARN] Test search failed: {e}")
        import traceback
        traceback.print_exc()

# =============================================================================
# Main ì‹¤í–‰
# =============================================================================
if __name__ == "__main__":
    print("=" * 80)
    print("PHASE 1: Document Indexing (txtai â†’ Milvus)")
    print("=" * 80)
    print("\nğŸ¯ í”Œë¡œìš°: txtai ë²¡í„°í™” â†’ Milvus ì €ì¥ì†Œ")
    print("=" * 80)
    
    try:
        # 1. txtai ì´ˆê¸°í™” (ë²¡í„° ìƒì„± ì—”ì§„)
        embeddings = initialize_embeddings()
        
        # 2. Milvus ì—°ê²° ë° ì»¬ë ‰ì…˜ ì„¤ì • (ì €ì¥ì†Œ)
        connect_milvus()
        collection = setup_milvus_collection(embed_dim=EMBED_DIM)
        
        # 3. ë¬¸ì„œ ë¡œë“œ
        documents = load_all_documents(DATA_DIR)
        
        if documents:
            # 4. âœ… txtaië¡œ ë²¡í„°í™” â†’ Milvusì— ì €ì¥
            vectorize_and_index_via_txtai(embeddings, collection, documents)
            
            # 5. í†µê³„
            show_collection_stats(collection)
            
            # 6. í…ŒìŠ¤íŠ¸ ê²€ìƒ‰ (txtaië¡œ ì¿¼ë¦¬ ë²¡í„°í™”)
            test_search(embeddings, collection, "ë”¥ëŸ¬ë‹", top_k=3)
        
        print("\n" + "=" * 80)
        print("âœ… Phase 1 ì™„ë£Œ: txtai ë²¡í„° â†’ Milvus ì €ì¥ë¨")
        print("=" * 80)
        print(f"\në‹¤ìŒ ë‹¨ê³„: python phase2_search_api.py ì‹¤í–‰")
    
    except Exception as e:
        print(f"\nâŒ ì—ëŸ¬ ë°œìƒ: {e}")
        import traceback
        traceback.print_exc()