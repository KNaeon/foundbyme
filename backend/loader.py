# loader.py
import os
from typing import Tuple

import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation
import markdown
import pytesseract
from PIL import Image

# Windows 환경에서 Tesseract 경로 설정 (필요시 주석 해제 및 경로 수정)
# pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 200) -> list[str]:
    if not text:
        return []
    
    # 너무 짧은 텍스트(예: 목차, 빈 페이지)는 무시
    if len(text.strip()) < 50:
        return []

    chunks = []
    start = 0
    text_len = len(text)
    
    while start < text_len:
        end = min(start + chunk_size, text_len)
        chunks.append(text[start:end])
        start += (chunk_size - overlap)
        
    return chunks

def load_text(path: str) -> list[dict]:
    """
    파일 경로 → [{'page': 1, 'content': '...'}, ...]
    """
    ext = os.path.splitext(path)[1].lstrip(".").lower()
    results = []

    if ext == "txt":
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            content = f.read()
        
        chunks = chunk_text(content)
        for i, chunk in enumerate(chunks):
            results.append({"page": i + 1, "content": chunk})

    elif ext == "pdf":
        try:
            with pdfplumber.open(path) as pdf:
                for i, page in enumerate(pdf.pages):
                    text = page.extract_text() or ""
                    if text.strip():
                        # PDF 페이지도 너무 길면 자를 수 있음 (선택 사항)
                        results.append({"page": i + 1, "content": text})
        except Exception as e:
            print(f"Error reading PDF {path}: {e}")

    elif ext == "docx":
        doc = DocxDocument(path)
        content = "\n".join(p.text for p in doc.paragraphs)
        
        chunks = chunk_text(content)
        for i, chunk in enumerate(chunks):
            results.append({"page": i + 1, "content": chunk})

    elif ext == "pptx":
        pres = Presentation(path)
        for i, slide in enumerate(pres.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
            content = "\n".join(texts)
            if content.strip():
                results.append({"page": i + 1, "content": content})

    elif ext == "md":
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            md_text = f.read()
        
        chunks = chunk_text(md_text)
        for i, chunk in enumerate(chunks):
            results.append({"page": i + 1, "content": chunk})

    elif ext in ["jpg", "jpeg", "png", "bmp", "tiff"]:
        try:
            image = Image.open(path)
            # 한국어+영어 추출
            text = pytesseract.image_to_string(image, lang='kor+eng')
            if text.strip():
                chunks = chunk_text(text)
                for i, chunk in enumerate(chunks):
                    results.append({"page": 1, "content": chunk})
        except Exception as e:
            print(f"Error reading Image {path}: {e}")
    
    return results
